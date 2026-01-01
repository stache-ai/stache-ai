"""PPTX (PowerPoint) document loader"""

from stache_ai.loaders.base import DocumentLoader


class PptxLoader(DocumentLoader):
    """Loader for Microsoft PowerPoint (.pptx) presentations"""

    @property
    def extensions(self) -> list[str]:
        return ['pptx']

    def load(self, file_path: str) -> str:
        # Lazy import
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

        return "\n\n".join(text_parts)
