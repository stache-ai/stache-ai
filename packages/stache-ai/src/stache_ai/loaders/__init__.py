"""Document loaders for various file formats"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def load_document(file_path: str, filename: str | None = None) -> str:
    """
    Load document and extract text

    Args:
        file_path: Path to the document file
        filename: Original filename (for extension detection)

    Returns:
        Extracted text content
    """
    path = Path(file_path)
    ext = (filename or path.name).lower().split('.')[-1]

    logger.info(f"Loading document: {filename or path.name} (type: {ext})")

    if ext in ['txt', 'md', 'markdown']:
        return load_text(file_path)
    elif ext == 'pdf':
        return load_pdf(file_path)
    elif ext == 'epub':
        return load_epub(file_path)
    elif ext == 'docx':
        return load_docx(file_path)
    elif ext == 'pptx':
        return load_pptx(file_path)
    elif ext == 'vtt':
        return load_vtt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def load_text(file_path: str) -> str:
    """Load plain text or markdown file"""
    with open(file_path, encoding='utf-8') as f:
        return f.read()


def load_pdf(file_path: str) -> str:
    """Load PDF file using pdfplumber, with OCR fallback for scanned PDFs"""
    import pdfplumber

    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)

    # If no text extracted, try OCR
    if not text_parts:
        logger.info(f"No text found in PDF, attempting OCR: {file_path}")
        text_parts = _ocr_pdf(file_path)

    return "\n\n".join(text_parts)


def _ocr_pdf(file_path: str) -> list[str]:
    """OCR a scanned PDF using ocrmypdf"""
    import subprocess
    import tempfile

    import pdfplumber

    text_parts = []
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp_path = tmp.name

        # Run ocrmypdf to add text layer
        result = subprocess.run(
            ['ocrmypdf', '--skip-text', '--quiet', file_path, tmp_path],
            capture_output=True,
            text=True
        )

        if result.returncode != 0:
            logger.warning(f"OCR failed: {result.stderr}")
            return []

        # Extract text from OCR'd PDF
        with pdfplumber.open(tmp_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

    except FileNotFoundError:
        logger.warning("ocrmypdf not installed, skipping OCR")
    except Exception as e:
        logger.warning(f"OCR error: {e}")
    finally:
        # Clean up temp file
        import os
        try:
            os.unlink(tmp_path)
        except:
            pass

    return text_parts


def load_epub(file_path: str) -> str:
    """Load EPUB file"""
    import ebooklib
    from bs4 import BeautifulSoup
    from ebooklib import epub

    book = epub.read_epub(file_path)
    text_parts = []

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text = soup.get_text()
            if text.strip():
                text_parts.append(text)

    return "\n\n".join(text_parts)


def load_docx(file_path: str) -> str:
    """Load Word document"""
    from docx import Document

    doc = Document(file_path)
    text_parts = [para.text for para in doc.paragraphs if para.text.strip()]

    return "\n\n".join(text_parts)


def load_pptx(file_path: str) -> str:
    """Load PowerPoint presentation"""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)

    return "\n\n".join(text_parts)


def load_vtt(file_path: str) -> str:
    """Load VTT (WebVTT) transcript file"""
    import webvtt

    text_parts = []
    for caption in webvtt.read(file_path):
        # Format: [timestamp] text
        text_parts.append(f"[{caption.start}] {caption.text}")

    return "\n\n".join(text_parts)
